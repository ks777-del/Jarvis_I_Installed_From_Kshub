from __future__ import annotations

from pathlib import Path
from typing import Dict, List

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class PPTBuilder:
    def __init__(self) -> None:
        self.base_dir = Path(__file__).resolve().parents[1]
        self.temp_dir = self.base_dir / "assets" / "cache"
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def build(
        self,
        brief: Dict[str, object],
        slides: List[Dict[str, object]],
        theme: Dict[str, object],
        assets: Dict[str, Dict[str, str]],
    ) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        for index, slide_data in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._paint_background(slide, theme, index)
            self._paint_shell(slide, theme)

            slide_id = f"slide_{index + 1}"
            if slide_data["type"] == "title":
                self._build_title(slide, slide_data, theme, assets.get(slide_id, {}))
            elif slide_data["type"] == "agenda":
                self._build_agenda(slide, slide_data, theme)
            elif slide_data["type"] == "visual":
                self._build_visual(slide, slide_data, theme, assets.get(slide_id, {}))
            elif slide_data["type"] == "data":
                self._build_data(slide, slide_data, theme)
            elif slide_data["type"] == "comparison":
                self._build_comparison(slide, slide_data, theme, assets.get(slide_id, {}))
            elif slide_data["type"] == "thank_you":
                self._build_thank_you(slide, slide_data, theme)
            else:
                self._build_content(slide, slide_data, theme)
        return prs

    def _paint_background(self, slide, theme: Dict[str, object], index: int) -> None:
        palette = theme["palette"]
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(palette["background"])
        bg.line.fill.background()

        glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.8), Inches(-0.8), Inches(4.2), Inches(4.2))
        glow.fill.solid()
        glow.fill.fore_color.rgb = self._rgb(palette["accent"])
        glow.fill.transparency = 0.78
        glow.line.fill.background()

        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.18))
        strip.fill.solid()
        strip.fill.fore_color.rgb = self._rgb(palette["secondary"] if index % 2 == 0 else palette["primary"])
        strip.line.fill.background()

    def _paint_shell(self, slide, theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.34), Inches(0.6), Inches(0.18), Inches(6.2))
        rail.fill.solid()
        rail.fill.fore_color.rgb = self._rgb(palette["primary"])
        rail.line.fill.background()

        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.4), Inches(0.55), Inches(1.15), Inches(0.38))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._rgb(palette["surface"])
        accent.line.fill.background()
        self._add_text(slide, theme, "JARVIS", (11.55, 0.58, 0.8, 0.25), 12, True, palette["primary"], align=PP_ALIGN.CENTER)

    def _build_title(self, slide, data: Dict[str, object], theme: Dict[str, object], asset: Dict[str, str]) -> None:
        palette = theme["palette"]
        hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.1), Inches(7.0), Inches(4.8))
        hero.fill.solid()
        hero.fill.fore_color.rgb = self._rgb(palette["surface"])
        hero.fill.transparency = 0.06
        hero.line.fill.background()

        self._add_text(slide, theme, data["title"], (1.2, 1.5, 6.2, 1.4), 34, True, palette["text_primary"])
        self._add_text(slide, theme, data["subtitle"], (1.2, 2.95, 5.8, 0.65), 19, False, palette["text_muted"])
        self._add_text(slide, theme, data["date"], (1.2, 4.85, 2.2, 0.4), 13, False, palette["secondary"])
        self._add_text(slide, theme, data["author"], (1.2, 5.25, 3.2, 0.4), 13, False, palette["text_muted"])

        image_path = asset.get("image")
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(image_path, Inches(8.55), Inches(1.05), width=Inches(3.8), height=Inches(5.2))
        else:
            panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(1.05), Inches(3.8), Inches(5.2))
            panel.fill.solid()
            panel.fill.fore_color.rgb = self._rgb(palette["surface"])
            panel.line.color.rgb = self._rgb(palette["accent"])

    def _build_agenda(self, slide, data: Dict[str, object], theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        self._add_title_block(slide, data["title"], theme)
        self._add_paragraph(slide, data.get("paragraph", ""), theme, (0.95, 1.45, 5.8, 0.55))
        for idx, bullet in enumerate(data.get("bullets", []), start=1):
            y = 2.15 + (idx - 1) * 0.78
            self._pill(slide, str(idx).zfill(2), (0.95, y, 0.55, 0.42), palette["primary"], theme)
            self._add_text(slide, theme, bullet, (1.65, y - 0.02, 4.9, 0.44), 17, True, palette["text_primary"])

        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(1.65), Inches(3.9), Inches(3.95))
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb(palette["surface"])
        card.line.fill.background()
        self._add_text(slide, theme, "Outcome", (8.45, 2.0, 2.8, 0.5), 24, True, palette["primary"])
        self._add_text(
            slide,
            theme,
            "A presentation that moves from orientation to insight and finishes with a confident takeaway.",
            (8.45, 2.75, 2.9, 1.6),
            16,
            False,
            palette["text_muted"],
        )

    def _build_content(self, slide, data: Dict[str, object], theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        self._add_title_block(slide, data["title"], theme)
        self._add_bullet_panel(slide, data["bullets"], theme, (0.95, 1.65, 6.0, 3.5))
        self._add_paragraph(slide, data.get("paragraph", ""), theme, (0.95, 5.45, 6.0, 0.85))

        quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.65), Inches(4.2), Inches(4.7))
        quote.fill.solid()
        quote.fill.fore_color.rgb = self._rgb(palette["surface"])
        quote.line.color.rgb = self._rgb(palette["accent"])
        self._add_text(slide, theme, "Strategic Lens", (8.35, 2.0, 2.9, 0.4), 22, True, palette["accent"])
        self._add_text(
            slide,
            theme,
            f"{data['title']} should feel practical, digestible, and ready for discussion.",
            (8.35, 2.85, 3.1, 1.2),
            16,
            False,
            palette["text_primary"],
        )
        self._add_text(
            slide,
            theme,
            "Keep the audience oriented with clean hierarchy and one clear idea per region.",
            (8.35, 4.3, 3.1, 1.1),
            15,
            False,
            palette["text_muted"],
        )

    def _build_visual(self, slide, data: Dict[str, object], theme: Dict[str, object], asset: Dict[str, str]) -> None:
        palette = theme["palette"]
        self._add_title_block(slide, data["title"], theme)
        image_path = asset.get("image")
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(image_path, Inches(0.95), Inches(1.65), width=Inches(7.25), height=Inches(4.85))
        else:
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.65), Inches(7.25), Inches(4.85))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb(palette["surface"])
            card.line.fill.background()
        self._add_text(slide, theme, data["caption"], (8.6, 2.05, 3.1, 0.9), 24, True, palette["text_primary"])
        self._add_text(slide, theme, data["visual_description"], (8.6, 3.0, 3.15, 1.5), 16, False, palette["text_muted"])
        self._pill(slide, "VISUAL", (8.62, 5.0, 1.2, 0.38), palette["primary"], theme)

    def _build_data(self, slide, data: Dict[str, object], theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        self._add_title_block(slide, data["title"], theme)
        chart_path = self._build_chart_image(data["metrics"], data.get("chart_type", "bar"), palette)
        slide.shapes.add_picture(str(chart_path), Inches(0.95), Inches(1.75), width=Inches(6.3), height=Inches(3.85))

        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.75), Inches(4.2), Inches(3.95))
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb(palette["surface"])
        card.line.fill.background()
        self._add_text(slide, theme, "Key Metrics", (8.35, 2.05, 2.9, 0.4), 22, True, palette["primary"])
        metric_lines = [f"{item['label']}: {item['value']}%" for item in data["metrics"]]
        self._add_bullets(slide, metric_lines, theme, (8.35, 2.75, 2.9, 1.8), palette["text_primary"])
        self._add_text(slide, theme, data["insight"], (8.35, 4.75, 3.0, 0.8), 14, False, palette["text_muted"])

    def _build_comparison(self, slide, data: Dict[str, object], theme: Dict[str, object], asset: Dict[str, str]) -> None:
        palette = theme["palette"]
        self._add_title_block(slide, data["title"], theme)
        self._comparison_panel(slide, data["left_title"], data["left_points"], theme, (0.95, 1.85, 5.2, 4.5), palette["secondary"])
        self._comparison_panel(slide, data["right_title"], data["right_points"], theme, (7.15, 1.85, 5.2, 4.5), palette["accent"])
        badge = asset.get("badge")
        if badge and Path(badge).exists():
            slide.shapes.add_picture(badge, Inches(6.17), Inches(3.0), width=Inches(1.0), height=Inches(1.0))

    def _build_thank_you(self, slide, data: Dict[str, object], theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(1.25), Inches(11.1), Inches(4.85))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(palette["surface"])
        panel.fill.transparency = 0.04
        panel.line.fill.background()
        self._add_text(slide, theme, data["title"], (3.2, 2.25, 6.8, 1.1), 36, True, palette["text_primary"], align=PP_ALIGN.CENTER)
        self._add_text(slide, theme, data["subtitle"], (3.0, 3.45, 7.0, 0.7), 19, False, palette["text_muted"], align=PP_ALIGN.CENTER)
        self._pill(slide, "JARVIS AI", (5.45, 4.65, 1.45, 0.42), palette["primary"], theme)

    def _add_title_block(self, slide, title: str, theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        self._add_text(slide, theme, title, (0.95, 0.72, 8.2, 0.8), 28, True, palette["text_primary"])
        self._add_text(slide, theme, "Designed and composed automatically by Jarvis.", (0.98, 1.18, 4.8, 0.3), 13, False, palette["text_muted"])

    def _add_bullet_panel(self, slide, bullets: List[str], theme: Dict[str, object], box) -> None:
        palette = theme["palette"]
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(palette["surface"])
        panel.fill.transparency = 0.03
        panel.line.fill.background()
        self._add_bullets(slide, bullets, theme, (box[0] + 0.35, box[1] + 0.35, box[2] - 0.7, box[3] - 0.7), palette["text_primary"])

    def _comparison_panel(self, slide, title: str, bullets: List[str], theme: Dict[str, object], box, line_color: str) -> None:
        palette = theme["palette"]
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(palette["surface"])
        panel.line.color.rgb = self._rgb(line_color)
        self._add_text(slide, theme, title, (box[0] + 0.3, box[1] + 0.25, box[2] - 0.6, 0.4), 21, True, line_color)
        self._add_bullets(slide, bullets, theme, (box[0] + 0.3, box[1] + 0.92, box[2] - 0.6, box[3] - 1.15), palette["text_primary"])

    def _add_paragraph(self, slide, text: str, theme: Dict[str, object], box) -> None:
        if not text:
            return
        self._add_text(slide, theme, text, box, 15, False, theme["palette"]["text_muted"])

    def _add_bullets(self, slide, bullets: List[str], theme: Dict[str, object], box, color: str) -> None:
        x, y, w, h = box
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.clear()
        for idx, bullet in enumerate(bullets):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = bullet
            paragraph.font.name = theme["fonts"]["body"]
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self._rgb(color)
            paragraph.space_after = Pt(10)
            paragraph.level = 0

    def _pill(self, slide, text: str, box, fill_color: str, theme: Dict[str, object]) -> None:
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
        pill.fill.solid()
        pill.fill.fore_color.rgb = self._rgb(fill_color)
        pill.line.fill.background()
        self._add_text(slide, theme, text, box, 11, True, "FFFFFF", align=PP_ALIGN.CENTER)

    def _add_text(self, slide, theme: Dict[str, object], text: str, box, size: int, bold: bool, color: str, align=PP_ALIGN.LEFT) -> None:
        x, y, w, h = box
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = text_box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.03)
        frame.margin_right = Inches(0.03)
        frame.margin_top = Inches(0.03)
        frame.margin_bottom = Inches(0.03)
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = theme["fonts"]["title"] if bold and size >= 22 else theme["fonts"]["body"]
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = self._rgb(color)
        paragraph.alignment = align

    def _build_chart_image(self, metrics: List[Dict[str, object]], chart_type: str, palette: Dict[str, str]) -> Path:
        path = self.temp_dir / f"chart_{chart_type}_{len(metrics)}.png"
        labels = [item["label"] for item in metrics]
        values = [item["value"] for item in metrics]
        fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=160)
        fig.patch.set_facecolor(f"#{palette['background']}")
        ax.set_facecolor(f"#{palette['surface']}")

        if chart_type == "pie":
            ax.pie(values, labels=labels, autopct="%1.0f%%", colors=[f"#{palette['primary']}", f"#{palette['secondary']}", f"#{palette['accent']}", "#94A3B8"])
        elif chart_type == "line":
            ax.plot(labels, values, color=f"#{palette['primary']}", linewidth=3, marker="o")
            ax.fill_between(labels, values, color=f"#{palette['accent']}", alpha=0.18)
        else:
            bars = ax.bar(labels, values, color=[f"#{palette['primary']}", f"#{palette['secondary']}", f"#{palette['accent']}", "#94A3B8"])
            for bar, value in zip(bars, values):
                ax.text(bar.get_x() + bar.get_width() / 2, value + 1.5, str(value), ha="center", fontsize=10, color=f"#{palette['text_primary']}")

        for spine in ("top", "right", "left", "bottom"):
            ax.spines[spine].set_visible(False)
        ax.tick_params(colors=f"#{palette['text_primary']}")
        ax.set_ylim(0, max(values) + 18)
        ax.grid(axis="y", linestyle="--", alpha=0.18)
        fig.tight_layout()
        fig.savefig(path, bbox_inches="tight", pad_inches=0.15)
        plt.close(fig)
        return path

    def _rgb(self, hex_color: str) -> RGBColor:
        value = hex_color.lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
