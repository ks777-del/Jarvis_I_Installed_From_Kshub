from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class ExportService:
    """Renders themed slide specs into a production-ready .pptx file."""

    def __init__(self) -> None:
        self.logger = logging.getLogger(__name__)

    def export(self, slide_specs: List[Dict[str, object]], theme: Dict[str, object], output_path: Path) -> None:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for spec in slide_specs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._paint_background(slide=slide, theme=theme)
            self._paint_decor(slide=slide, theme=theme)
            self._render_slide(slide=slide, spec=spec, theme=theme)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        self.logger.info("Presentation exported to %s", output_path)

    def _paint_background(self, slide, theme: Dict[str, object]) -> None:
        palette = theme["palette"]
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(palette["background"])
        bg.line.fill.background()

    def _paint_decor(self, slide, theme: Dict[str, object]) -> None:
        palette = theme["palette"]

        left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
        left.fill.solid()
        left.fill.fore_color.rgb = self._rgb(palette["primary"])
        left.line.fill.background()

        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.22), Inches(0), Inches(13.11), Inches(0.14))
        top.fill.solid()
        top.fill.fore_color.rgb = self._rgb(palette["secondary"])
        top.line.fill.background()

        orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.8), Inches(0.35), Inches(1.0), Inches(1.0))
        orb.fill.solid()
        orb.fill.fore_color.rgb = self._rgb(palette["accent"])
        orb.fill.transparency = 0.2
        orb.line.fill.background()

    def _render_slide(self, slide, spec: Dict[str, object], theme: Dict[str, object]) -> None:
        slide_type = spec["slide_type"]
        if slide_type == "title":
            self._render_title_slide(slide, spec, theme)
        elif slide_type == "comparison":
            self._render_comparison_slide(slide, spec, theme)
        elif slide_type == "highlight":
            self._render_highlight_slide(slide, spec, theme)
        elif slide_type == "bullet":
            self._render_bullet_slide(slide, spec, theme)
        else:
            self._render_content_slide(slide, spec, theme)

    def _render_title_slide(self, slide, spec: Dict[str, object], theme: Dict[str, object]) -> None:
        layout = spec["layout"]
        self._add_text(
            slide,
            text=spec["title"],
            box=layout["title_box"],
            font_name=theme["fonts"]["title"],
            size=theme["font_sizes"]["title"],
            bold=True,
            color=theme["palette"]["text_primary"],
        )
        self._add_text(
            slide,
            text=spec.get("subtitle", ""),
            box=layout["subtitle_box"],
            font_name=theme["fonts"]["subtitle"],
            size=theme["font_sizes"]["subtitle"],
            bold=False,
            color=theme["palette"]["text_muted"],
        )

    def _render_content_slide(self, slide, spec: Dict[str, object], theme: Dict[str, object]) -> None:
        layout = spec["layout"]
        self._add_section_title(slide, spec["title"], layout["title_box"], theme)
        self._add_bullets(slide, spec.get("bullets", []), layout["body_box"], theme)

    def _render_bullet_slide(self, slide, spec: Dict[str, object], theme: Dict[str, object]) -> None:
        self._render_content_slide(slide, spec, theme)

    def _render_highlight_slide(self, slide, spec: Dict[str, object], theme: Dict[str, object]) -> None:
        layout = spec["layout"]
        palette = theme["palette"]

        self._add_section_title(slide, spec["title"], layout["title_box"], theme)

        hx, hy, hw, hh = layout["highlight_box"]
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(hx), Inches(hy), Inches(hw), Inches(hh))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(palette["surface"])
        panel.line.color.rgb = self._rgb(palette["secondary"])

        self._add_text(
            slide,
            text=spec.get("highlight", ""),
            box=layout["highlight_box"],
            font_name=theme["fonts"]["highlight"],
            size=theme["font_sizes"]["highlight"],
            bold=False,
            color=palette["text_primary"],
            margin=0.15,
        )
        self._add_bullets(slide, spec.get("bullets", []), layout["body_box"], theme)

    def _render_comparison_slide(self, slide, spec: Dict[str, object], theme: Dict[str, object]) -> None:
        layout = spec["layout"]
        palette = theme["palette"]

        self._add_section_title(slide, spec["title"], layout["title_box"], theme)

        self._add_panel_with_bullets(
            slide,
            title="Traditional",
            bullets=spec.get("left_items", []),
            box=layout["left_box"],
            title_color=palette["primary"],
            theme=theme,
        )
        self._add_panel_with_bullets(
            slide,
            title="Modern",
            bullets=spec.get("right_items", []),
            box=layout["right_box"],
            title_color=palette["accent"],
            theme=theme,
        )

    def _add_section_title(self, slide, title: str, box, theme: Dict[str, object]) -> None:
        self._add_text(
            slide,
            text=title,
            box=box,
            font_name=theme["fonts"]["title"],
            size=theme["font_sizes"]["section_title"],
            bold=True,
            color=theme["palette"]["text_primary"],
        )

    def _add_bullets(self, slide, bullets: List[str], box, theme: Dict[str, object]) -> None:
        if not bullets:
            return

        x, y, w, h = box
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = text_box.text_frame
        tf.clear()

        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.name = theme["fonts"]["body"]
            p.font.size = Pt(theme["font_sizes"]["bullet"])
            p.font.color.rgb = self._rgb(theme["palette"]["text_primary"])
            p.space_after = Pt(12)
            p.alignment = PP_ALIGN.LEFT

    def _add_panel_with_bullets(self, slide, title: str, bullets: List[str], box, title_color: str, theme: Dict[str, object]) -> None:
        x, y, w, h = box
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(theme["palette"]["surface"])
        panel.line.color.rgb = self._rgb(title_color)

        self._add_text(
            slide,
            text=title,
            box=(x + 0.2, y + 0.15, w - 0.4, 0.55),
            font_name=theme["fonts"]["title"],
            size=theme["font_sizes"]["subtitle"],
            bold=True,
            color=title_color,
        )
        self._add_bullets(slide, bullets, (x + 0.2, y + 0.9, w - 0.4, h - 1.1), theme)

    def _add_text(self, slide, text: str, box, font_name: str, size: int, bold: bool, color: str, margin: float = 0.05) -> None:
        x, y, w, h = box
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = text_box.text_frame
        tf.margin_left = Inches(margin)
        tf.margin_right = Inches(margin)
        tf.margin_top = Inches(margin)
        tf.margin_bottom = Inches(margin)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = self._rgb(color)
        p.alignment = PP_ALIGN.LEFT

    def _rgb(self, hex_color: str) -> RGBColor:
        value = hex_color.lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
